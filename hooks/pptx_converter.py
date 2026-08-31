# -*- coding: utf-8 -*-
"""
MkDocs build hook: convert .pptx files in docs/ to browsable slide pages.

Priority:
  1. Slide images (full layout fidelity) via LibreOffice or PowerPoint COM.
  2. Fallback: text/table/image extraction via python-pptx.

Usage:
  1. Drop .pptx files anywhere under docs/
  2. Run mkdocs build or mkdocs serve
  3. Each .pptx gets a same-named .md with a slide viewer.
"""

import hashlib
import shutil
import subprocess
import urllib.parse
import uuid
from pathlib import Path


# ---------------------------------------------------------------------------
# Slide image conversion
# ---------------------------------------------------------------------------

def _convert_with_libreoffice(pptx_path: Path, out_dir: Path) -> list[Path] | None:
    """Convert PPTX to PNGs using LibreOffice + pdftoppm. Returns image paths."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    pdftoppm = shutil.which("pdftoppm")
    if not pdftoppm:
        return None

    tmp = out_dir / ".convert"
    tmp.mkdir(parents=True, exist_ok=True)

    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf",
             "--outdir", str(tmp), str(pptx_path)],
            check=True, capture_output=True, timeout=180,
        )
        pdf_path = tmp / (pptx_path.stem + ".pdf")
        if not pdf_path.exists():
            return None

        prefix = out_dir / "slide"
        subprocess.run(
            [pdftoppm, "-png", "-r", "160", str(pdf_path), str(prefix)],
            check=True, capture_output=True, timeout=180,
        )

        images = sorted(out_dir.glob("slide-*.png"))
        if not images:
            return None

        # Normalize to slide-01.png ... slide-NN.png
        normalized = []
        for i, img in enumerate(images, 1):
            target = out_dir / f"slide-{i:02d}.png"
            if img != target:
                img.replace(target)
            normalized.append(target)
        return normalized
    except Exception:
        return None
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def _convert_with_powerpoint(pptx_path: Path, out_dir: Path) -> list[Path] | None:
    """Convert PPTX to PNGs using PowerPoint COM via PowerShell (Windows)."""
    if Path("/").anchor != "\\":
        return None  # not Windows

    ps = r"""
$ErrorActionPreference = "Stop"
$ppt = New-Object -ComObject PowerPoint.Application
try {{
  $pres = $ppt.Presentations.Open('{pptx}', $true, $false, $false)
  $pres.Export('{outdir}', 'PNG', 1600, 900)
  $pres.Close()
}} finally {{
  $ppt.Quit()
  [System.Runtime.InteropServices.Marshal]::ReleaseComObject($ppt) | Out-Null
}}
""".format(pptx=str(pptx_path.resolve()).replace("'", "''"),
           outdir=str(out_dir.resolve()).replace("'", "''"))

    script = out_dir / "_pptx_export.ps1"
    script.write_bytes(ps.encode("utf-8-sig"))
    try:
        subprocess.run(
            ["powershell", "-NoProfile", "-ExecutionPolicy", "Bypass",
             "-File", str(script)],
            check=True, capture_output=True, timeout=300,
        )
    except Exception:
        return None
    finally:
        script.unlink(missing_ok=True)

    # PowerPoint exports localized names such as 幻灯片1.PNG; find newest batch.
    pngs = sorted(
        (p for p in out_dir.glob("*.png")
         if not p.name.startswith("slide-")),
        key=lambda p: p.stat().st_mtime,
    )
    if not pngs:
        return None

    normalized = []
    for i, p in enumerate(pngs, 1):
        target = out_dir / f"slide-{i:02d}.png"
        p.replace(target)
        normalized.append(target)
    return normalized


def _get_slide_images(pptx_path: Path, out_dir: Path) -> list[Path] | None:
    """Return fresh slide images, converting first if needed."""
    out_dir.mkdir(parents=True, exist_ok=True)
    images = sorted(out_dir.glob("slide-*.png"))
    marker = out_dir / ".source-hash"
    pptx_hash = hashlib.md5(pptx_path.read_bytes()).hexdigest()

    # Reuse committed images only when the PPTX source hash still matches.
    if (images and marker.exists()
            and marker.read_text(encoding="utf-8").strip() == pptx_hash):
        return images

    for old in images:
        old.unlink(missing_ok=True)
    marker.unlink(missing_ok=True)

    converted = (_convert_with_libreoffice(pptx_path, out_dir)
                 or _convert_with_powerpoint(pptx_path, out_dir))
    if converted:
        marker.write_text(pptx_hash, encoding="utf-8")
    return converted


# ---------------------------------------------------------------------------
# Fallback text extraction (python-pptx)
# ---------------------------------------------------------------------------

def _extract_pptx_slides(pptx_path, img_output_dir):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(pptx_path))
    slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {"number": slide_num, "title": "", "elements": []}
        _first_text = None

        if slide.shapes.title and slide.shapes.title.has_text_frame:
            t = slide.shapes.title.text.strip()
            if t:
                slide_data["title"] = t

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue

            elem = None
            if shape.has_text_frame and not shape.has_table:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        p = {"text": text, "level": para.level or 0}
                        try:
                            p["bold"] = bool(para.runs[0].font.bold) if para.runs else False
                        except Exception:
                            p["bold"] = False
                        paragraphs.append(p)
                if paragraphs:
                    if _first_text is None:
                        _first_text = paragraphs[0]["text"]
                    elem = {"type": "text", "paragraphs": paragraphs}

            if shape.has_table:
                rows_data = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                if rows_data:
                    elem = {"type": "table", "rows": rows_data}

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image = shape.image
                    ext_map = {"image/png": "png", "image/jpeg": "jpg",
                               "image/gif": "gif", "image/bmp": "bmp"}
                    ext = ext_map.get(image.content_type, "png")
                    img_id = f"s{slide_num}_{uuid.uuid4().hex[:8]}"
                    img_path = img_output_dir / f"{img_id}.{ext}"
                    img_path.parent.mkdir(parents=True, exist_ok=True)
                    img_path.write_bytes(image.blob)
                    elem = {"type": "image",
                            "src": f"../assets/pptx/{img_id}.{ext}",
                            "alt": f"Slide {slide_num} image"}
                except Exception:
                    pass

            if elem:
                slide_data["elements"].append(elem)

        if not slide_data["title"] and _first_text:
            slide_data["title"] = _first_text

        slides.append(slide_data)

    return slides


# ---------------------------------------------------------------------------
# Markdown renderers
# ---------------------------------------------------------------------------

def _relative_img_url(img: Path, docs_dir: Path, md_parent: Path) -> str:
    """Return a MkDocs-friendly relative image URL for a markdown page."""
    rel = img.relative_to(docs_dir).as_posix()
    # The generated page lives in its own directory (md file name + "/index.html"),
    # so we need one extra ".." on top of md_parent's depth.
    depth = len(md_parent.relative_to(docs_dir).parts) + 1
    return "../" * depth + rel


def _render_slideshow_markdown(images: list[Path], title: str,
                               docs_dir: Path, md_parent: Path,
                               version: str = "",
                               interactive_url: str | None = None) -> str:
    """Render slide images as an interactive viewer."""
    cache_suffix = f"?v={version}" if version else ""
    imgs = []
    for p in sorted(images, key=lambda x: x.name):
        url = _relative_img_url(p, docs_dir, md_parent)
        imgs.append(url + cache_suffix)
    img_json = ",\n".join(f'    "{u}"' for u in imgs)
    interactive_section = ""
    if interactive_url:
        interactive_section = f"""
## 互动演示

<div class="ppt-interactive">
  <iframe src="{interactive_url}" width="100%" height="600px" frameborder="0"
          allowfullscreen title="PPT 互动演示"></iframe>
</div>

<style>
.ppt-interactive {{ max-width: 1000px; margin: 0 auto 24px; }}
.ppt-interactive iframe {{ width: 100%; border: 0; border-radius: 8px; background: #fff; }}
</style>
"""

    return f"""# {title}

{interactive_section}

!!! info "共 {len(imgs)} 页幻灯片"
    此页面保留了 PPT 原始布局，可使用按钮或键盘左右方向键翻页。

<div class="slide-viewer">
  <div class="slide-stage">
    <img id="slide-main" src="{imgs[0]}" alt="{title}">
    <button class="slide-nav prev" onclick="goSlide(-1)" aria-label="上一页">&#10094;</button>
    <button class="slide-nav next" onclick="goSlide(1)" aria-label="下一页">&#10095;</button>
    <div class="slide-counter" id="slide-counter">1 / {len(imgs)}</div>
    <button class="slide-play" id="slide-play" onclick="togglePlay()" aria-label="播放">
      <span class="play-icon">&#9654;</span>
    </button>
  </div>
  <div class="slide-thumbs" id="slide-thumbs"></div>
</div>

<style>
.slide-viewer {{ max-width: 1000px; margin: 0 auto; }}
.slide-stage {{
  position: relative; background: #1a1b26; border-radius: 8px;
  box-shadow: 0 12px 32px rgba(0,0,0,.28); overflow: hidden;
}}
.slide-stage img {{ width: 100%; height: auto; display: block; }}
.slide-nav {{
  position: absolute; top: 50%; transform: translateY(-50%);
  width: 42px; height: 42px; border: 0; border-radius: 50%;
  background: rgba(15,17,26,.72); color: #fff; font-size: 18px;
  cursor: pointer; display: flex; align-items: center; justify-content: center;
}}
.slide-nav:hover {{ background: rgba(94,108,255,.85); }}
.slide-nav.prev {{ left: 10px; }}
.slide-nav.next {{ right: 10px; }}
.slide-play {{
  position: absolute; bottom: 10px; left: 12px;
  width: 40px; height: 40px; border: 0; border-radius: 50%;
  background: rgba(15,17,26,.72); color: #fff; font-size: 16px;
  cursor: pointer; display: flex; align-items: center; justify-content: center;
}}
.slide-play:hover {{ background: rgba(94,108,255,.85); }}
.slide-counter {{
  position: absolute; bottom: 10px; right: 12px;
  background: rgba(15,17,26,.78); color: #e5e7eb;
  padding: 3px 10px; border-radius: 999px; font-size: 12px;
}}
.slide-thumbs {{
  display: flex; gap: 8px; overflow-x: auto;
  margin-top: 12px; padding-bottom: 4px;
}}
.slide-thumbs button {{
  flex: 0 0 92px; padding: 0; border: 2px solid transparent;
  border-radius: 6px; overflow: hidden; cursor: pointer; background: none;
}}
.slide-thumbs button.active {{ border-color: #5e6cff; }}
.slide-thumbs img {{ width: 100%; height: 56px; object-fit: cover; display: block; }}
</style>

<script>
var slideImages = [
{img_json}
];
var slideIndex = 0;
var isPlaying = false;
var playTimer = null;
var mainImg = document.getElementById("slide-main");
var counter = document.getElementById("slide-counter");
var thumbs = document.getElementById("slide-thumbs");
var playBtn = document.getElementById("slide-play");
var playIcon = playBtn.querySelector(".play-icon");

function renderThumbs() {{
  thumbs.innerHTML = "";
  slideImages.forEach(function(src, i) {{
    var btn = document.createElement("button");
    btn.innerHTML = '<img src="' + src + '" alt="slide ' + (i+1) + '">';
    if (i === slideIndex) btn.className = "active";
    btn.onclick = function() {{ slideIndex = i; update(); }};
    thumbs.appendChild(btn);
  }});
}}

function update() {{
  mainImg.src = slideImages[slideIndex];
  counter.textContent = (slideIndex + 1) + " / " + slideImages.length;
  var btns = thumbs.children;
  for (var i = 0; i < btns.length; i++) {{
    btns[i].className = (i === slideIndex) ? "active" : "";
  }}
}}

function goSlide(delta) {{
  slideIndex = (slideIndex + delta + slideImages.length) % slideImages.length;
  update();
}}

function togglePlay() {{
  isPlaying = !isPlaying;
  if (isPlaying) {{
    playIcon.innerHTML = "&#10074;&#10074;";
    playTimer = setInterval(function() {{
      slideIndex = (slideIndex + 1) % slideImages.length;
      update();
    }}, 3000);
  }} else {{
    playIcon.innerHTML = "&#9654;";
    clearInterval(playTimer);
    playTimer = null;
  }}
}}

document.addEventListener("keydown", function(e) {{
  if (e.key === "ArrowLeft") goSlide(-1);
  if (e.key === "ArrowRight") goSlide(1);
  if (e.key === " ") {{
    e.preventDefault();
    togglePlay();
  }}
}});

renderThumbs();
update();
</script>
"""


def _render_text_markdown(slides: list, title: str) -> str:
    """Fallback renderer: extracted text content."""
    lines = [f"# {title}", "",
             f'!!! warning "当前为文本提取模式"',
             "    未找到 PowerPoint 或 LibreOffice，无法生成保真幻灯片图片，仅展示提取的文字。", ""]

    for s in slides:
        slide_title = s["title"] or f"第 {s['number']} 页"
        lines.append(f"## {slide_title}")
        lines.append("")
        if not s["elements"]:
            lines.append("*（此页无可提取的文本内容）*")
            lines.append("")
            continue
        for elem in s["elements"]:
            if elem["type"] == "text":
                for p in elem["paragraphs"]:
                    text = p["text"].replace("\n", " ")
                    if p.get("bold"):
                        text = f"**{text}**"
                    lines.append("    " * p.get("level", 0) + f"- {text}")
                lines.append("")
            elif elem["type"] == "table":
                rows = elem["rows"]
                if not rows:
                    continue
                lines.append("| " + " | ".join(str(c) for c in rows[0]) + " |")
                lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
                for row in rows[1:]:
                    padded = list(row) + [""] * (len(rows[0]) - len(row))
                    lines.append("| " + " | ".join(str(c) for c in padded[:len(rows[0])]) + " |")
                lines.append("")
            elif elem["type"] == "image":
                lines.append(f'<figure markdown="span">')
                lines.append(f'  ![{elem.get("alt", "image")}]({elem["src"]})')
                lines.append(f"</figure>")
                lines.append("")
        lines.append("---")
        lines.append("")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# MkDocs hook entry point
# ---------------------------------------------------------------------------

def on_pre_build(*, config, **kwargs):
    docs_dir = Path(config["docs_dir"])
    assets_img_dir = docs_dir / "assets" / "pptx"
    assets_img_dir.mkdir(parents=True, exist_ok=True)

    pptx_files = list(docs_dir.rglob("*.pptx"))
    if not pptx_files:
        print("[pptx-hook] No PPTX files found, skipping.")
        return

    print(f"[pptx-hook] Found {len(pptx_files)} PPTX file(s), converting...")

    for pptx_path in pptx_files:
        rel = pptx_path.relative_to(docs_dir)
        md_path = pptx_path.with_suffix(".md")
        md_parent = md_path.parent
        title = pptx_path.stem

        print(f"  -> {rel}")

        images = _get_slide_images(pptx_path, assets_img_dir)
        if images:
            version = hashlib.md5(pptx_path.read_bytes()).hexdigest()[:10]
            interactive_url = None
            site_url = config.get("site_url") or ""
            if site_url:
                public_pptx = site_url.rstrip("/") + "/" + rel.as_posix()
                encoded_src = urllib.parse.quote(public_pptx, safe="")
                interactive_url = (
                    "https://view.officeapps.live.com/op/embed.aspx?src="
                    + encoded_src
                )
            markdown = _render_slideshow_markdown(
                images, title, docs_dir, md_parent, version, interactive_url)
            print(f"     [OK] Generated {len(images)} slide images")
        else:
            slides = _extract_pptx_slides(pptx_path, assets_img_dir)
            if not slides:
                print("     [WARN] No content extracted")
                continue
            markdown = _render_text_markdown(slides, title)
            print(f"     [OK] Text fallback ({len(slides)} slides)")

        md_path.write_text(markdown, encoding="utf-8")

    print("[pptx-hook] Done.")
