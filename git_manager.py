import os
import git
from pathlib import Path
from datetime import datetime


class GitManager:
    """Manages git operations for a watched repository."""
    
    def __init__(self, repo_path: str):
        self.repo_path = Path(repo_path).resolve()
        self._repo = None
    
    @property
    def repo(self):
        if self._repo is None:
            self._repo = self._init_or_open()
        return self._repo
    
    def _init_or_open(self):
        """Initialize a new git repo or open an existing one."""
        self.repo_path.mkdir(parents=True, exist_ok=True)
        try:
            r = git.Repo(self.repo_path)
            if r.head.is_valid():
                return r
            return self._create_first_commit(r)
        except git.InvalidGitRepositoryError:
            r = git.Repo.init(self.repo_path)
            return self._create_first_commit(r)
    
    def _create_first_commit(self, r):
        """Create initial commit so HEAD is valid."""
        readme = self.repo_path / "README.md"
        if not readme.exists():
            readme.write_text(f"# {self.repo_path.name}\n\nAuto-watched repository.\n")
        r.index.add(["README.md"])
        r.index.commit("Initial commit")
        return r
    
    def auto_commit(self, message=None):
        """Stage all changes and commit. Returns True if a commit was made."""
        r = self.repo
        if not r.is_dirty(untracked_files=True):
            return False
        r.git.add(A=True)
        if message is None:
            changed = [i.a_path for i in r.index.diff(None)]
            untracked = r.untracked_files
            files = changed + untracked
            msg = f"Auto-sync: {', '.join(files[:5])}"
            if len(files) > 5:
                msg += f" (+{len(files) - 5} more)"
            message = msg
        r.index.commit(message)
        return True
    
    def get_commits(self, max_count=50):
        """Return list of recent commits."""
        out = []
        for c in self.repo.iter_commits(max_count=max_count):
            out.append({
                "sha": c.hexsha[:8],
                "message": c.message.strip(),
                "author": str(c.author),
                "date": datetime.fromtimestamp(c.committed_date).isoformat(),
                "date_str": datetime.fromtimestamp(c.committed_date).strftime("%Y-%m-%d %H:%M"),
            })
        return out
    
    def get_file_tree(self, base_path=None):
        """Build a recursive file-tree list from the filesystem."""
        base = self.repo_path if base_path is None else Path(base_path)
        items = []
        for child in sorted(base.iterdir(), key=lambda x: (x.is_file(), x.name.lower())):
            name = child.name
            if name == ".git" or name.startswith("."):
                continue
            rel = str(child.relative_to(self.repo_path)).replace("\\", "/")
            if child.is_dir():
                items.append({
                    "name": name, "path": rel, "type": "directory",
                    "children": self.get_file_tree(child),
                })
            else:
                items.append({"name": name, "path": rel, "type": "file"})
        return items
    
    def get_file_content(self, file_path, commit_sha=None):
        """Read file content (UTF-8), optionally at a past commit."""
        if commit_sha:
            try:
                return self.repo.git.show(f"{commit_sha}:{file_path}")
            except git.exc.GitCommandError:
                return None
        full = self.repo_path / file_path
        if full.exists() and full.is_file():
            try:
                return full.read_text(encoding="utf-8")
            except (UnicodeDecodeError, OSError):
                return None
        return None
    

    def get_pptx_content(self, file_path):
        """Extract all visual elements from a PPTX file (text, tables, images)."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            import uuid, shutil
            
            full_path = self.repo_path / file_path
            if not full_path.exists():
                return None
            prs = Presentation(str(full_path))
            
            # Clean and create image directory for this PPTX
            img_dir = self.repo_path / ".pptx_images"
            if img_dir.exists():
                shutil.rmtree(img_dir)
            img_dir.mkdir(parents=True)
            
            repo_name = self.repo_path.name
            slides = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_data = {
                    "number": slide_num,
                    "title": "",
                    "elements": []
                }
                # Title
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    t = slide.shapes.title.text.strip()
                    if t:
                        slide_data["title"] = t
                
                for shape in slide.shapes:
                    if shape == slide.shapes.title:
                        continue
                    
                    elem = None
                    
                    # --- Text element ---
                    if shape.has_text_frame:
                        paragraphs = []
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                p = {"text": text, "level": para.level}
                                try:
                                    p["bold"] = bool(para.font.bold)
                                except:
                                    p["bold"] = False
                                try:
                                    if para.font.size:
                                        p["font_size"] = para.font.size.pt
                                except:
                                    pass
                                paragraphs.append(p)
                        if paragraphs:
                            elem = {"type": "text", "paragraphs": paragraphs}
                    
                    # --- Table element ---
                    if shape.has_table:
                        table = shape.table
                        rows_data = []
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            rows_data.append(cells)
                        if rows_data:
                            elem = {
                                "type": "table",
                                "rows": rows_data,
                                "num_rows": len(table.rows),
                                "num_cols": len(table.columns),
                            }
                    
                    # --- Image element ---
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            ext_map = {"image/png": "png", "image/jpeg": "jpg",
                                       "image/gif": "gif", "image/bmp": "bmp"}
                            ext = ext_map.get(image.content_type, "png")
                            img_id = f"s{slide_num}_{uuid.uuid4().hex[:8]}"
                            img_path = img_dir / f"{img_id}.{ext}"
                            with open(img_path, "wb") as f:
                                f.write(image.blob)
                            elem = {
                                "type": "image",
                                "id": img_id,
                                "ext": ext,
                                "img_url": f"/repos/{repo_name}/raw/.pptx_images/{img_id}.{ext}",
                            }
                        except Exception:
                            pass
                    
                    if elem:
                        slide_data["elements"].append(elem)
                
                slides.append(slide_data)
            
            return slides
        except ImportError:
            return None
        except Exception as e:
            import traceback
            traceback.print_exc()
            return None

    def get_docx_content(self, file_path):
        """Extract elements from a Word document (paragraphs, tables)."""
        try:
            from docx import Document

            full_path = self.repo_path / file_path
            if not full_path.exists():
                return None
            doc = Document(str(full_path))
            body = doc.element.body
            elements = []

            # Collect only paragraphs that are direct children of body (skip table-internal ones)
            body_paras = {}
            for para in doc.paragraphs:
                parent = para._element.getparent()
                if parent is not None and parent.tag.split("}")[-1] != "body":
                    continue
                text = para.text.strip()
                if not text:
                    continue

                style_name = para.style.name if para.style else "Normal"
                etype = "paragraph"
                hlevel = 0
                if style_name.startswith("Heading"):
                    try:
                        hlevel = int(style_name.split()[-1])
                        etype = "heading"
                    except:
                        pass
                elif "List Bullet" in style_name:
                    etype = "list_bullet"
                elif "List Number" in style_name:
                    etype = "list_number"

                elem = {"type": etype, "text": text, "level": hlevel}
                for run in para.runs:
                    if run.bold:
                        elem["bold"] = True
                    if run.italic:
                        elem["italic"] = True
                body_paras[para._element] = elem

            # Collect tables
            body_tables = {}
            for table in doc.tables:
                rows_data = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows_data.append(cells)
                body_tables[table._element] = {"type": "table", "rows": rows_data}

            # Iterate body children in document order
            for child in body:
                tag = child.tag.split("}")[-1]
                if tag == "p" and child in body_paras:
                    elements.append(body_paras[child])
                elif tag == "tbl" and child in body_tables:
                    elements.append(body_tables[child])

            return elements
        except ImportError:
            return None
        except Exception as e:
            import traceback
            traceback.print_exc()
            return None


    def get_repo_info(self):
        """Return summary info for this repo."""
        r = self.repo
        commits = self.get_commits(1)
        return {
            "name": self.repo_path.name,
            "path": str(self.repo_path),
            "branch": r.active_branch.name,
            "status": "dirty" if r.is_dirty(untracked_files=True) else "clean",
            "file_count": sum(1 for f in self.repo_path.rglob("*") if f.is_file() and ".git" not in str(f)),
            "last_commit": commits[0] if commits else None,
        }
